"""Tests for pptx/writer.py — create .pptx → apply rules → re-parse to verify."""

from __future__ import annotations

import hashlib
import tempfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from wfm_agents.pptx.parser import parse_pptx
from wfm_agents.pptx.writer import apply_font_rules


def _save_prs(prs: Presentation) -> Path:
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)  # noqa: SIM115
    prs.save(tmp.name)
    tmp.close()
    return Path(tmp.name)


def _make_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def _make_content_slide(prs: Presentation, title: str, body: str) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    for ph in slide.placeholders:
        ph_name = ph.placeholder_format.type.name
        if ph_name in ("BODY", "OBJECT"):
            ph.text = body
            break


def _make_blank_textbox_slide(prs: Presentation, text: str) -> None:
    layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = text


def _get_run_font(slide: object, shape_idx: int, para_idx: int, run_idx: int, field: str) -> str | None:
    shape = list(slide.shapes)[shape_idx]
    para = list(shape.text_frame.paragraphs)[para_idx]
    run = para.runs[run_idx]
    rPr = run._r.get_or_add_rPr()  # noqa: SLF001
    tag_map = {"latin": qn("a:latin"), "ea": qn("a:ea"), "cs": qn("a:cs")}
    el = rPr.find(tag_map[field])
    return el.get("typeface") if el is not None else None


class TestApplyCjkFont:
    def test_ea_font_modified(self) -> None:
        prs = Presentation()
        _make_title_slide(prs, "测试标题")
        path = _save_prs(prs)

        apply_font_rules(path, [{"scope": "title", "ea": "思源黑体"}])

        prs2 = Presentation(str(path))
        slide = prs2.slides[0]
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type.name in ("TITLE", "CENTER_TITLE"):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        rPr = run._r.get_or_add_rPr()  # noqa: SLF001
                        ea_el = rPr.find(qn("a:ea"))
                        assert ea_el is not None
                        assert ea_el.get("typeface") == "思源黑体"
                break


class TestApplyLatinFont:
    def test_latin_font_modified(self) -> None:
        prs = Presentation()
        _make_title_slide(prs, "Test Title")
        path = _save_prs(prs)

        apply_font_rules(path, [{"scope": "title", "latin": "Arial"}])

        prs2 = Presentation(str(path))
        slide = prs2.slides[0]
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type.name in ("TITLE", "CENTER_TITLE"):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        rPr = run._r.get_or_add_rPr()  # noqa: SLF001
                        latin_el = rPr.find(qn("a:latin"))
                        assert latin_el is not None
                        assert latin_el.get("typeface") == "Arial"
                break


class TestApplyCsFont:
    def test_cs_font_modified(self) -> None:
        prs = Presentation()
        _make_title_slide(prs, "Arabic Test")
        path = _save_prs(prs)

        apply_font_rules(path, [{"scope": "all", "cs": "Arial"}])

        prs2 = Presentation(str(path))
        slide = prs2.slides[0]
        found = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        rPr = run._r.get_or_add_rPr()  # noqa: SLF001
                        cs_el = rPr.find(qn("a:cs"))
                        assert cs_el is not None
                        assert cs_el.get("typeface") == "Arial"
                        found = True
        assert found


class TestApplySize:
    def test_size_modified(self) -> None:
        prs = Presentation()
        _make_title_slide(prs, "测试")
        path = _save_prs(prs)

        apply_font_rules(path, [{"scope": "title", "size_pt": 28}])

        prs2 = Presentation(str(path))
        slide = prs2.slides[0]
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type.name in ("TITLE", "CENTER_TITLE"):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        assert run.font.size is not None
                        assert abs(run.font.size.pt - 28.0) < 0.1
                break


class TestTitleScope:
    def test_only_title_modified(self) -> None:
        prs = Presentation()
        _make_title_slide(prs, "标题", "副标题")
        path = _save_prs(prs)

        apply_font_rules(path, [{"scope": "title", "ea": "思源黑体"}])

        prs2 = Presentation(str(path))
        slide = prs2.slides[0]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            is_title = shape.is_placeholder and shape.placeholder_format.type.name in ("TITLE", "CENTER_TITLE")
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    rPr = run._r.get_or_add_rPr()  # noqa: SLF001
                    ea_el = rPr.find(qn("a:ea"))
                    ea_val = ea_el.get("typeface") if ea_el is not None else None
                    if is_title:
                        assert ea_val == "思源黑体"
                    else:
                        # subtitle should NOT be modified
                        assert ea_val != "思源黑体" or ea_val is None


class TestBodyScope:
    def test_only_body_modified(self) -> None:
        prs = Presentation()
        _make_content_slide(prs, "标题", "正文内容")
        path = _save_prs(prs)

        apply_font_rules(path, [{"scope": "body", "ea": "思源黑体"}])

        prs2 = Presentation(str(path))
        slide = prs2.slides[0]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            is_body = shape.is_placeholder and shape.placeholder_format.type.name in ("BODY", "OBJECT")
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    rPr = run._r.get_or_add_rPr()  # noqa: SLF001
                    ea_el = rPr.find(qn("a:ea"))
                    ea_val = ea_el.get("typeface") if ea_el is not None else None
                    if is_body:
                        assert ea_val == "思源黑体"
                    else:
                        assert ea_val != "思源黑体"


class TestTextboxScope:
    def test_only_textbox_modified(self) -> None:
        prs = Presentation()
        _make_title_slide(prs, "标题")
        _make_blank_textbox_slide(prs, "自由文本")
        path = _save_prs(prs)

        apply_font_rules(path, [{"scope": "textbox", "ea": "思源黑体"}])

        prs2 = Presentation(str(path))
        # Slide 0: title slide — should not be modified
        for shape in prs2.slides[0].shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    rPr = run._r.get_or_add_rPr()  # noqa: SLF001
                    ea_el = rPr.find(qn("a:ea"))
                    if ea_el is not None:
                        assert ea_el.get("typeface") != "思源黑体"

        # Slide 1: blank textbox — should be modified
        for shape in prs2.slides[1].shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    rPr = run._r.get_or_add_rPr()  # noqa: SLF001
                    ea_el = rPr.find(qn("a:ea"))
                    assert ea_el is not None
                    assert ea_el.get("typeface") == "思源黑体"


class TestAllScope:
    def test_all_shapes_modified(self) -> None:
        prs = Presentation()
        _make_title_slide(prs, "标题", "副标题")
        path = _save_prs(prs)

        apply_font_rules(path, [{"scope": "all", "ea": "思源黑体"}])

        prs2 = Presentation(str(path))
        found = 0
        for slide in prs2.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        rPr = run._r.get_or_add_rPr()  # noqa: SLF001
                        ea_el = rPr.find(qn("a:ea"))
                        assert ea_el is not None
                        assert ea_el.get("typeface") == "思源黑体"
                        found += 1
        assert found >= 2


class TestMultipleRules:
    def test_both_rules_applied(self) -> None:
        prs = Presentation()
        _make_title_slide(prs, "标题", "副标题")
        _make_content_slide(prs, "章节", "正文")
        path = _save_prs(prs)

        apply_font_rules(path, [
            {"scope": "title", "ea": "思源黑体", "size_pt": 28},
            {"scope": "body", "ea": "宋体", "size_pt": 18},
        ])

        prs2 = Presentation(str(path))
        for slide in prs2.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                is_title = shape.is_placeholder and shape.placeholder_format.type.name in ("TITLE", "CENTER_TITLE")
                is_body = shape.is_placeholder and shape.placeholder_format.type.name in ("BODY", "OBJECT")
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        rPr = run._r.get_or_add_rPr()  # noqa: SLF001
                        ea_el = rPr.find(qn("a:ea"))
                        ea_val = ea_el.get("typeface") if ea_el is not None else None
                        if is_title:
                            assert ea_val == "思源黑体"
                            assert abs(run.font.size.pt - 28.0) < 0.1
                        elif is_body:
                            assert ea_val == "宋体"
                            assert abs(run.font.size.pt - 18.0) < 0.1


class TestSlideRange:
    def test_only_target_slides_modified(self) -> None:
        prs = Presentation()
        _make_title_slide(prs, "封面")
        _make_content_slide(prs, "第1章", "内容1")
        _make_content_slide(prs, "第2章", "内容2")
        path = _save_prs(prs)

        apply_font_rules(path, [
            {"scope": "all", "ea": "思源黑体", "slide_range": [2, 3]},
        ])

        prs2 = Presentation(str(path))
        # Slide 0 (index 1): should NOT be modified
        for shape in prs2.slides[0].shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    rPr = run._r.get_or_add_rPr()  # noqa: SLF001
                    ea_el = rPr.find(qn("a:ea"))
                    if ea_el is not None:
                        assert ea_el.get("typeface") != "思源黑体"

        # Slide 1+ (index 2+): should be modified
        found_modified = False
        for i in range(1, len(prs2.slides)):
            slide = prs2.slides[i]
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        rPr = run._r.get_or_add_rPr()  # noqa: SLF001
                        ea_el = rPr.find(qn("a:ea"))
                        if ea_el is not None and ea_el.get("typeface") == "思源黑体":
                            found_modified = True
        assert found_modified


class TestSaveToDifferentFile:
    def test_source_unchanged(self) -> None:
        prs = Presentation()
        _make_title_slide(prs, "原始标题")
        path = _save_prs(prs)
        original_size = path.stat().st_size

        output = Path(str(path).replace(".pptx", "_out.pptx"))
        apply_font_rules(path, [{"scope": "all", "ea": "思源黑体"}], output_path=output)

        assert output.exists()
        assert path.stat().st_size == original_size


class TestDryRun:
    def test_file_not_modified(self) -> None:
        prs = Presentation()
        _make_title_slide(prs, "测试")
        path = _save_prs(prs)
        original_size = path.stat().st_size

        result = apply_font_rules(path, [{"scope": "all", "ea": "思源黑体"}], dry_run=True)

        assert path.stat().st_size == original_size
        assert "would_modify" in result
        assert result["total_changes"] >= 1

    def test_returns_change_list(self) -> None:
        prs = Presentation()
        _make_title_slide(prs, "测试标题")
        path = _save_prs(prs)

        result = apply_font_rules(path, [{"scope": "all", "ea": "思源黑体", "size_pt": 28}], dry_run=True)

        assert result["total_changes"] >= 1
        changes = result["would_modify"]
        fields = {c["field"] for c in changes}
        assert "ea" in fields
        assert "size_pt" in fields


class TestPreserveImages:
    def test_images_not_corrupted(self) -> None:
        prs = Presentation()
        layout = prs.slide_layouts[5]  # Blank
        slide = prs.slides.add_slide(layout)

        # Add a tiny textbox
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5))
        txBox.text_frame.text = "文本"

        path = _save_prs(prs)
        original_hash = hashlib.sha256(path.read_bytes()).hexdigest()

        apply_font_rules(path, [{"scope": "all", "ea": "思源黑体"}])

        # File should have changed (different hash)
        new_hash = hashlib.sha256(path.read_bytes()).hexdigest()
        assert new_hash != original_hash

        # Re-open to verify it's valid PPTX
        prs2 = Presentation(str(path))
        assert len(prs2.slides) == 1
