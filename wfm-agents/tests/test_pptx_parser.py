"""Tests for pptx/parser.py — use python-pptx to build fixture presentations."""

from __future__ import annotations

import tempfile
import unittest.mock
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from wfm_agents.pptx.parser import (
    format_pptx_content,
    parse_pptx,
    summarize_fonts,
)
from wfm_agents.pptx.parser import _INHERITED


def _save_prs(prs: Presentation) -> Path:
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)  # noqa: SIM115
    prs.save(tmp.name)
    tmp.close()
    return Path(tmp.name)


def _make_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def _make_content_slide(prs: Presentation, title: str, body: str) -> None:
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    for ph in slide.placeholders:
        ph_name = ph.placeholder_format.type.name
        if ph_name in ("BODY", "OBJECT"):
            ph.text = body
            break


class TestParsePptx:
    def test_empty_presentation(self) -> None:
        prs = Presentation()
        path = _save_prs(prs)
        result = parse_pptx(path)
        assert result["metadata"]["slide_count"] == 0
        assert result["slides"] == []
        assert result["stats"]["shapes_total"] == 0

    def test_title_slide(self) -> None:
        prs = Presentation()
        _make_title_slide(prs, "项目汇报", "2024年度总结")
        path = _save_prs(prs)

        result = parse_pptx(path)
        assert len(result["slides"]) == 1

        slide = result["slides"][0]
        title_shapes = [s for s in slide["shapes"] if s["placeholder_type"] == "title"]
        assert len(title_shapes) >= 1
        title_texts = [r["text"] for s in title_shapes for p in s["paragraphs"] for r in p["runs"]]
        assert "项目汇报" in title_texts

    def test_content_slide_body(self) -> None:
        prs = Presentation()
        _make_content_slide(prs, "第一章", "这是正文内容")
        path = _save_prs(prs)

        result = parse_pptx(path)
        slide = result["slides"][0]
        body_shapes = [s for s in slide["shapes"] if s["placeholder_type"] == "body"]
        assert len(body_shapes) >= 1
        body_texts = [r["text"] for s in body_shapes for p in s["paragraphs"] for r in p["runs"]]
        assert "这是正文内容" in body_texts

    def test_font_extraction(self) -> None:
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        title = slide.shapes.title
        run = title.text_frame.paragraphs[0].add_run()
        run.text = "测试"
        run.font.name = "Arial"
        run.font.size = Pt(32)

        from pptx.oxml.ns import qn
        from lxml import etree

        rPr = run._r.get_or_add_rPr()  # noqa: SLF001
        ea_el = etree.SubElement(rPr, qn("a:ea"))
        ea_el.set("typeface", "微软雅黑")

        path = _save_prs(prs)
        result = parse_pptx(path)

        slide_data = result["slides"][0]
        font = slide_data["shapes"][0]["paragraphs"][0]["runs"][0]["font"]
        assert font["latin"] == "Arial"
        assert font["ea"] == "微软雅黑"
        assert font["size_pt"] == 32.0

    def test_inherited_font(self) -> None:
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        title = slide.shapes.title
        run = title.text_frame.paragraphs[0].add_run()
        run.text = "无字体设置"

        path = _save_prs(prs)
        result = parse_pptx(path)

        slide_data = result["slides"][0]
        font = slide_data["shapes"][0]["paragraphs"][0]["runs"][0]["font"]
        # latin should be resolved from theme or marked inherited
        assert font["ea"] is not None

    def test_textbox_shape(self) -> None:
        prs = Presentation()
        layout = prs.slide_layouts[5]  # Blank
        slide = prs.slides.add_slide(layout)

        from pptx.util import Inches

        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tf = txBox.text_frame
        tf.text = "自由文本框"

        path = _save_prs(prs)
        result = parse_pptx(path)

        slide_data = result["slides"][0]
        textbox_shapes = [s for s in slide_data["shapes"] if s["shape_type"] == "textbox"]
        assert len(textbox_shapes) >= 1
        assert textbox_shapes[0]["placeholder_type"] is None

    def test_table_shape(self) -> None:
        prs = Presentation()
        layout = prs.slide_layouts[5]  # Blank
        slide = prs.slides.add_slide(layout)

        rows, cols = 2, 2
        table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(4), Inches(2))
        table = table_shape.table
        table.cell(0, 0).text = "名称"
        table.cell(0, 1).text = "金额"
        table.cell(1, 0).text = "钢材"
        table.cell(1, 1).text = "5000"

        path = _save_prs(prs)
        result = parse_pptx(path)

        slide_data = result["slides"][0]
        table_shapes = [s for s in slide_data["shapes"] if s["shape_type"] == "table"]
        assert len(table_shapes) >= 1

        all_texts = [r["text"] for p in table_shapes[0]["paragraphs"] for r in p["runs"]]
        assert "名称" in all_texts
        assert "5000" in all_texts

    def test_multiple_slides(self) -> None:
        prs = Presentation()
        _make_title_slide(prs, "标题1")
        _make_content_slide(prs, "标题2", "内容2")
        path = _save_prs(prs)

        result = parse_pptx(path)
        assert len(result["slides"]) == 2
        assert result["slides"][0]["index"] == 1
        assert result["slides"][1]["index"] == 2
        assert result["stats"]["slides_total"] == 2

    def test_file_size_limit(self) -> None:
        prs = Presentation()
        _make_title_slide(prs, "test")
        path = _save_prs(prs)

        with unittest.mock.patch.object(Path, "stat") as mock_stat:
            mock_stat.return_value.st_size = 25 * 1024 * 1024
            with pytest.raises(ValueError, match="文件过大"):
                parse_pptx(path)

    def test_max_slides_truncation(self) -> None:
        prs = Presentation()
        for i in range(5):
            _make_title_slide(prs, f"Slide {i}")
        path = _save_prs(prs)

        result = parse_pptx(path, max_slides=3)
        assert len(result["slides"]) == 3
        assert result["stats"]["truncated"] is True
        assert result["stats"]["slides_total"] == 5

    def test_metadata(self) -> None:
        prs = Presentation()
        prs.core_properties.title = "测试标题"
        prs.core_properties.author = "测试作者"
        _make_title_slide(prs, "test")
        path = _save_prs(prs)

        result = parse_pptx(path)
        assert result["metadata"]["title"] == "测试标题"
        assert result["metadata"]["author"] == "测试作者"
        assert result["metadata"]["slide_width_pt"] > 0
        assert result["metadata"]["slide_height_pt"] > 0


class TestSummarizeFonts:
    def test_basic_summary(self) -> None:
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        run = slide.shapes.title.text_frame.paragraphs[0].add_run()
        run.text = "测试"
        run.font.name = "Arial"
        run.font.size = Pt(24)

        path = _save_prs(prs)
        content = parse_pptx(path)
        summary = summarize_fonts(content)

        assert summary["total_runs"] >= 1
        assert "Arial" in summary["fonts"]["latin"]
        assert 24.0 in summary["sizes"]["unique"]

    def test_multiple_fonts(self) -> None:
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        from pptx.oxml.ns import qn
        from lxml import etree

        run1 = slide.shapes.title.text_frame.paragraphs[0].add_run()
        run1.text = "英文"
        run1.font.name = "Calibri"

        run2 = slide.shapes.title.text_frame.paragraphs[0].add_run()
        run2.text = "中文"
        run2.font.name = "Arial"
        rPr = run2._r.get_or_add_rPr()  # noqa: SLF001
        ea_el = etree.SubElement(rPr, qn("a:ea"))
        ea_el.set("typeface", "微软雅黑")

        path = _save_prs(prs)
        content = parse_pptx(path)
        summary = summarize_fonts(content)

        assert "Calibri" in summary["fonts"]["latin"]
        assert "微软雅黑" in summary["fonts"]["ea"]


class TestFormatPptxContent:
    def test_basic_formatting(self) -> None:
        prs = Presentation()
        _make_title_slide(prs, "项目汇报", "2024年度总结")
        path = _save_prs(prs)

        content = parse_pptx(path)
        md = format_pptx_content(content)

        assert "# Presentation:" in md
        assert "Slide 1" in md
        assert "项目汇报" in md

    def test_truncation_hint(self) -> None:
        content = {
            "metadata": {"title": "大型演示", "slide_count": 100,
                         "slide_width_pt": 960, "slide_height_pt": 540},
            "slides": [{"index": i, "layout": "Title", "shapes": []} for i in range(1, 4)],
            "stats": {"slides_total": 100, "slides_returned": 3,
                      "truncated": False, "shapes_total": 0, "runs_total": 0},
        }
        md = format_pptx_content(content, max_slides_preview=2)
        assert "剩余" in md
        assert "100 张" in md
